import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 16:9 Widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color Palette (Industrial Modern / Cyber IoT)
BG_DARK = RGBColor(15, 23, 42)       # Slate 900
CARD_BG = RGBColor(30, 41, 59)       # Slate 800
CARD_BORDER = RGBColor(51, 65, 85)   # Slate 700
ACCENT_BLUE = RGBColor(14, 165, 233) # Sky 500
ACCENT_CYAN = RGBColor(6, 182, 212)  # Cyan 500
ACCENT_EMERALD = RGBColor(16, 185, 129) # Emerald 500
ACCENT_AMBER = RGBColor(245, 158, 11)   # Amber 500
TEXT_WHITE = RGBColor(248, 250, 252)    # Slate 50
TEXT_MUTED = RGBColor(148, 163, 184)    # Slate 400
TEXT_DARK = RGBColor(15, 23, 42)

def set_slide_background(slide, color=BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, slide_num, title_text, category="IoT MOTOR MONITORING & PREDICTIVE MAINTENANCE"):
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p0 = tf.paragraphs[0]
    p0.text = f"SLIDE {slide_num:02d}  |  {category.upper()}"
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_CYAN
    p0.font.name = "Segoe UI"
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.font.name = "Segoe UI"

def create_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_slide_layout = prs.slide_layouts[6]
    
    # SLIDE 1: TITLE SLIDE
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s1)
    
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(3.2), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(14, 116, 144)
    badge.line.fill.background()
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "INDUSTRY 4.0 IoT PLATFORM"
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = RGBColor(224, 242, 254)
    p_b.alignment = PP_ALIGN.CENTER
    
    tb_title = s1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(2.2))
    tf_t = tb_title.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "IoT-Based Motor Monitoring and\nPredictive Maintenance System"
    p_t.font.size = Pt(36)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE
    p_t.font.name = "Segoe UI"
    
    create_card(s1, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.4))
    tb_p = s1.shapes.add_textbox(Inches(1.1), Inches(4.5), Inches(5.0), Inches(2.0))
    tf_p = tb_p.text_frame
    p = tf_p.paragraphs[0]
    p.text = "PROJECT PRESENTERS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    
    presenters = [
        "Person 1: Hardware & ESP32 Embedded Systems",
        "Person 2: Web Dashboard & FastAPI Backend Server",
        "Person 3: Machine Learning & Predictive Analytics"
    ]
    for pres in presenters:
        p = tf_p.add_paragraph()
        p.text = f"•  {pres}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
    
    create_card(s1, Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.4))
    tb_tech = s1.shapes.add_textbox(Inches(7.1), Inches(4.5), Inches(5.1), Inches(2.0))
    tf_tech = tb_tech.text_frame
    p = tf_tech.paragraphs[0]
    p.text = "CORE TECHNOLOGIES"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    
    techs = [
        "Hardware: ESP32 Microcontroller | Multi-Sensor Bus (I2C/ADC)",
        "Sensors: DHT22 (Temp/Hum) | ACS712 (Current) | MPU6050 (Vibe)",
        "Backend: Python | FastAPI | SQLite | WebSocket Server",
        "Frontend & ML: React Dashboard | TailwindCSS | Scikit-Learn"
    ]
    for tech in techs:
        p = tf_tech.add_paragraph()
        p.text = f"✔  {tech}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
        
    add_speaker_notes(s1, "Welcome to our presentation on the IoT-Based Motor Monitoring and Predictive Maintenance System. Our team combines embedded hardware, cloud/backend software, and machine learning.")

    # SLIDE 2: INTRODUCTION
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s2)
    add_header(s2, 1, "Introduction — Industrial Motors & Maintenance")
    
    create_card(s2, Inches(0.8), Inches(1.7), Inches(5.7), Inches(4.0))
    tb = s2.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.1), Inches(3.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Role of Industrial Motors"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    
    bullets = [
        "Essential driving forces in manufacturing, pumps, conveyors, compressors, and cooling fans.",
        "Continuous 24/7 operation under demanding mechanical and thermal loads.",
        "Unexpected failures lead to severe downtime and economic losses."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        
    p = tf.add_paragraph()
    p.text = "\nConsequences of Unplanned Failure:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AMBER
    for risk in ["Production downtime", "Increased repair costs", "Equipment damage", "Workplace safety hazards"]:
        p = tf.add_paragraph()
        p.text = f" ✖ {risk}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    create_card(s2, Inches(6.8), Inches(1.7), Inches(5.7), Inches(4.0))
    tb = s2.shapes.add_textbox(Inches(7.1), Inches(1.9), Inches(5.1), Inches(3.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Our IoT-Based Approach"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    
    sol_bullets = [
        "Continuous telemetry data collection using multi-sensor nodes connected to an ESP32 microcontroller.",
        "High-speed wireless Wi-Fi transmission to a centralized FastAPI server.",
        "Live industrial web dashboard for real-time health visualization and automated alarms.",
        "Machine Learning pipeline for early anomaly detection and predictive failure modeling."
    ]
    for b in sol_bullets:
        p = tf.add_paragraph()
        p.text = f"✔ {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE

    create_card(s2, Inches(0.8), Inches(5.9), Inches(11.733), Inches(1.1), bg_color=RGBColor(24, 39, 75))
    tb = s2.shapes.add_textbox(Inches(1.1), Inches(6.0), Inches(11.1), Inches(0.9))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary Quote:"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = '“Our project continuously monitors industrial motor condition via sensors. The ESP32 collects telemetry over Wi-Fi, the dashboard provides live operator insights, and ML predicts impending faults.”'
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = RGBColor(224, 242, 254)
    add_speaker_notes(s2, "Industrial motors power factories and pumps. Unplanned downtime is costly. We propose an end-to-end IoT and ML platform.")

    # SLIDE 3: ABSTRACT
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s3)
    add_header(s3, 2, "Project Abstract")
    
    create_card(s3, Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.9))
    tb = s3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "01. EDGE SENSING"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    items = [
        "ESP32 Microcontroller: 32-bit dual core CPU, integrated 2.4GHz Wi-Fi.",
        "DHT22: High-accuracy ambient & motor surface temp/humidity.",
        "ACS712: Hall-effect current transducer with calibrated zero offset.",
        "MPU6050: 3-Axis accelerometer for dynamic vibration spectrum analysis.",
        "IR Optical Sensor: Object / RPM state detection."
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
        
    create_card(s3, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.9))
    tb = s3.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "02. CLOUD & BACKEND"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    items = [
        "FastAPI Server: High-throughput asynchronous REST API & WebSocket pipeline.",
        "SQLite Persistence: Time-series historical logging of all motor metrics.",
        "Real-Time Telemetry: Low-latency streaming of sensor data to operators.",
        "Command & Control: Remote speed adjustment and emergency shutdown capabilities."
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    create_card(s3, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.9))
    tb = s3.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "03. PREDICTIVE AI"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    items = [
        "Condition Assessment: Automated vibration & current anomaly grading.",
        "ML Modeling: Trend analysis over temperature rise, vibration spikes, and current draw.",
        "Goal: Shift from reactive/calendar maintenance to condition-based predictive maintenance.",
        "Result: Reduced downtime, lower maintenance expenditure, improved equipment lifespan."
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s3, "This abstract outlines the complete pipeline: hardware edge sensing, cloud backend ingestion, and predictive ML intelligence.")

    # SLIDE 4: SECTOR & APPLICATION AREAS
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s4)
    add_header(s4, 3, "Target Sector & Suitable Application Areas")
    
    apps = [
        ("Manufacturing Industries", "Production assembly lines, automated conveyors, and precision robotic machinery requiring continuous uptime.", ACCENT_CYAN),
        ("Pumping & Water Systems", "Municipal water pumps, chemical delivery pumps, and industrial cooling circulation loops.", ACCENT_BLUE),
        ("HVAC & Ventilation", "Large commercial air handlers, industrial exhaust fans, and cooling tower motor drives.", ACCENT_EMERALD),
        ("Mining & Heavy Industry", "Heavy rock crushers, ore conveyor belts, vibratory screens, and underground ventilation.", ACCENT_AMBER),
        ("Process Industries", "Chemical agitators, food processing mixers, packaging lines, and textile spinning mills.", RGBColor(168, 85, 247)),
        ("Smart Factory / Industry 4.0", "Fully connected cyber-physical factory environments with centralized telemetry and predictive maintenance.", RGBColor(236, 72, 153))
    ]
    for idx, (title, desc, col) in enumerate(apps):
        row = idx // 3
        col_idx = idx % 3
        left = Inches(0.8 + col_idx * 4.0)
        top = Inches(1.8 + row * 2.5)
        create_card(s4, left, top, Inches(3.7), Inches(2.2))
        
        tb = s4.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(3.3), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{idx+1}. {title}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s4, "Target sector is industrial manufacturing, especially continuous-duty operations where unexpected motor stoppage incurs massive costs.")

    # SLIDE 5: PROBLEM STATEMENT
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s5)
    add_header(s5, 4, "Problem Statement — Traditional vs Modern Maintenance")
    
    create_card(s5, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s5.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Traditional Maintenance Limitations"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AMBER
    trad = [
        "Periodic Manual Inspections: Only captures state at discrete, infrequent intervals.",
        "Blind Spots: Developing faults between scheduled inspections go undetected.",
        "Undetected Overheating: Gradual thermal breakdown of winding insulation.",
        "Undetected Vibration: Mechanical imbalance or bearing degradation worsens over time.",
        "High Labor & Time: Requires dedicated technicians to manually visit each machine."
    ]
    for item in trad:
        p = tf.add_paragraph()
        p.text = f"✖ {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    create_card(s5, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s5.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Proposed IoT Condition-Based Solution"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    sol = [
        "Continuous 24/7 Telemetry: High-frequency sampling of key physical variables.",
        "Real-Time Anomaly Alerts: Immediate notifications when vibration, current, or temp exceed thresholds.",
        "Automated Edge Calibration: Calibrated sensor zero-point offsets ensure precision.",
        "Centralized Web Monitoring: Single-pane-of-glass dashboard for plant-wide visibility.",
        "Machine Learning Forecasting: Early warning before catastrophic mechanical failure."
    ]
    for item in sol:
        p = tf.add_paragraph()
        p.text = f"✔ {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s5, "Traditional maintenance is periodic and blind to sudden faults. Our IoT solution provides 24/7 continuous health tracking.")

    # SLIDE 6: OBJECTIVES
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s6)
    add_header(s6, 5, "System Objectives")
    
    create_card(s6, Inches(0.8), Inches(1.8), Inches(11.733), Inches(1.3), bg_color=RGBColor(24, 39, 75))
    tb = s6.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(11.1), Inches(1.1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "PRIMARY OBJECTIVE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = "To engineer an end-to-end real-time IoT platform that continuously monitors industrial motor physical parameters, stores historical telemetry, provides interactive web visualization, and establishes the foundation for predictive maintenance ML models."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE

    specs = [
        ("Thermal & Ambient Sensing", "Track motor housing temperature and ambient humidity with DHT22 to avoid overheating.", ACCENT_CYAN),
        ("Current & Load Tracking", "Measure instantaneous current with ACS712 to detect overload and resistance.", ACCENT_BLUE),
        ("Vibration Spectrum Analysis", "Compute 3-axis acceleration and vibration deviation with MPU6050 accelerometer.", ACCENT_EMERALD),
        ("Motor Driver & Control", "Interface L298N dual H-Bridge driver with ESP32 for PWM speed control and safety shutoff.", ACCENT_AMBER),
        ("Backend & Web Dashboard", "Stream data via Wi-Fi to FastAPI server, store in SQLite, and render in live UI.", RGBColor(168, 85, 247)),
        ("Predictive Intelligence", "Prepare structured time-series datasets for machine learning failure prediction.", RGBColor(236, 72, 153))
    ]
    for idx, (title, desc, col) in enumerate(specs):
        row = idx // 3
        col_idx = idx % 3
        left = Inches(0.8 + col_idx * 4.0)
        top = Inches(3.3 + row * 1.8)
        create_card(s6, left, top, Inches(3.7), Inches(1.6))
        
        tb = s6.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(3.3), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s6, "Our objectives span edge sensing, motor control, backend telemetry pipeline, and predictive AI data modeling.")

    # SLIDE 7: HARDWARE REQUIREMENTS
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s7)
    add_header(s7, 6, "Hardware Requirements & Bill of Materials")
    
    rows, cols = 11, 2
    table_shape = s7.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
    table = table_shape.table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(8.233)
    
    hw_data = [
        ("Component", "Functional Purpose in System"),
        ("ESP32 Dev Board", "Central 32-bit dual-core microcontroller with 2.4 GHz Wi-Fi / Bluetooth"),
        ("DHT22 Sensor", "Digital high-precision temperature & humidity sensing"),
        ("ACS712 (30A/20A/5A)", "Hall-effect isolated motor current transducer (calibrated at 1.8925 V zero-point)"),
        ("MPU6050 IMU", "3-Axis acceleration & vibration magnitude measurement via I2C bus"),
        ("IR Optical Sensor", "Digital infrared sensor for object detection / RPM gating"),
        ("L298N Motor Driver", "Dual H-Bridge power driver for DC motor speed (PWM) and direction control"),
        ("DC Motor", "Industrial motor prototype under continuous health monitoring"),
        ("External Power Supply", "Regulated 12V/5V DC power source for motor and sensor rails"),
        ("Resistor Divider Network", "Signal level shifting and ADC protection for ESP32 3.3V inputs"),
        ("Breadboard & Jumper Wires", "Physical prototyping harness and electrical interconnects")
    ]
    for r_idx, (c1, c2) in enumerate(hw_data):
        cell1 = table.cell(r_idx, 0)
        cell2 = table.cell(r_idx, 1)
        cell1.text = c1
        cell2.text = c2
        for c, text_val in [(cell1, c1), (cell2, c2)]:
            c.fill.solid()
            if r_idx == 0:
                c.fill.fore_color.rgb = RGBColor(14, 116, 144)
            else:
                c.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(24, 32, 47)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(12 if r_idx == 0 else 11)
                p.font.bold = (r_idx == 0 or c == cell1)
                p.font.color.rgb = TEXT_WHITE if r_idx == 0 else (ACCENT_CYAN if c == cell1 else TEXT_WHITE)
                p.font.name = "Segoe UI"
    add_speaker_notes(s7, "Comprehensive BOM covering the ESP32, sensor array, L298N driver, power supply, and protection circuitry.")

    # SLIDE 8: SOFTWARE REQUIREMENTS
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s8)
    add_header(s8, 7, "Software Architecture & Tech Stack")
    
    sw_stacks = [
        ("Embedded Firmware", ["Arduino IDE / C++", "ESP32 Board Core (v2.x/v3.x)", "Wire.h (I2C Bus Protocol)", "Adafruit DHT & Sensor Libs", "Non-blocking timer loops"], ACCENT_CYAN),
        ("Backend Server", ["Python 3.11+", "FastAPI Web Framework", "Uvicorn ASGI Server", "Pydantic Schema Validation", "SQLite Time-Series DB"], ACCENT_BLUE),
        ("Web Frontend Dashboard", ["Modern React 19 + TypeScript", "TailwindCSS Industrial Design", "Recharts Interactive Charts", "WebSocket Real-Time Sync", "Responsive Operator UI"], ACCENT_EMERALD),
        ("Machine Learning & Comms", ["Scikit-learn / NumPy / Pandas", "Vibration FFT Feature Extraction", "Anomaly Detection Classifiers", "HTTP REST JSON Ingestion", "WebSocket Broadcast Pipeline"], ACCENT_AMBER)
    ]
    for idx, (title, items, col) in enumerate(sw_stacks):
        left = Inches(0.8 + idx * 3.0)
        top = Inches(1.8)
        create_card(s8, left, top, Inches(2.75), Inches(5.0))
        tb = s8.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), Inches(2.45), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        for it in items:
            p = tf.add_paragraph()
            p.text = f"✔ {it}"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s8, "Software architecture spanning embedded C++ on ESP32, Python FastAPI backend, React dashboard, and ML libraries.")

    # SLIDE 9: SENSOR EXPLANATION (ESP32 & DHT22)
    s9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s9)
    add_header(s9, 8, "Core Controller & Thermal Sensor")
    
    create_card(s9, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s9.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ESP32 Microcontroller"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    esp_features = [
        "Core: Xtensa Dual-Core 32-bit LX6 microprocessor operating up to 240 MHz.",
        "Wireless: Integrated 802.11 b/g/n Wi-Fi + BLE 4.2 for direct cloud connectivity.",
        "Analog Peripherals: Multi-channel 12-bit SAR ADC for analog sensor acquisition.",
        "Digital Bus: Hardware I2C (GPIO 21/22) for high-speed MPU6050 accelerometer reads.",
        "Role: Edge data acquisition, physical unit conversion, and HTTP JSON transmission."
    ]
    for it in esp_features:
        p = tf.add_paragraph()
        p.text = f"• {it}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    create_card(s9, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s9.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DHT22 Temperature & Humidity"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    dht_features = [
        "Temperature Range: -40°C to +80°C (±0.5°C accuracy).",
        "Humidity Range: 0 to 100% RH (±2% accuracy).",
        "Protocol: Single-bus digital signal on ESP32 GPIO 4.",
        "Industrial Significance: Motor winding insulation degrades exponentially when operating above rated thermal class.",
        "Early Detection: Detects thermal runaway and inadequate ventilation before winding burnout."
    ]
    for it in dht_features:
        p = tf.add_paragraph()
        p.text = f"• {it}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s9, "ESP32 provides high processing power and Wi-Fi. DHT22 tracks motor temperature to prevent insulation degradation.")

    # SLIDE 10: ACS712 CURRENT SENSOR
    s10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s10)
    add_header(s10, 9, "ACS712 Motor Current Sensing & Calibration")
    
    create_card(s10, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s10.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Operating Principle & Significance"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    acs_pts = [
        "Hall-Effect Transducer: Provides galvanic isolation between high-power motor line and ESP32.",
        "Linear Output: Yields an analog voltage proportional to AC/DC load current.",
        "Current Signatures:",
        "  • Inrush Current: Spike during motor startup.",
        "  • Steady-State Load: Nominal operating current.",
        "  • Overload / Jamming: Sustained current surge indicating mechanical binding or stalled rotor."
    ]
    for it in acs_pts:
        p = tf.add_paragraph()
        p.text = it
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    create_card(s10, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s10.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Mathematical Model & Calibration"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    calib_pts = [
        "Measured Zero-Current Voltage: 1.8925 V (Calibrated Offset)",
        "Sensor Sensitivity (ACS712-20A/30A): ~100 mV/A or 66 mV/A",
        "Conversion Formula in ESP32:",
        "  V_sensor = (ADC_raw / 4095.0) * V_ref * Divider_Ratio",
        "  I_motor = (V_sensor - V_zero_offset) / Sensitivity",
        "Accuracy Optimization: 100-sample moving average filter eliminates switching noise and ADC jitter."
    ]
    for it in calib_pts:
        p = tf.add_paragraph()
        p.text = it
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s10, "ACS712 measures current via Hall effect. Our hardware calibration identified a 1.8925V zero-point offset for accurate Amperage calculation.")

    # SLIDE 11: MPU6050 VIBRATION SENSOR
    s11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s11)
    add_header(s11, 10, "MPU6050 Vibration & Acceleration Analysis")
    
    create_card(s11, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s11.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Accelerometer Vector Mathematics"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    vibe_pts = [
        "MPU6050 6-DOF IMU: High-precision 3-axis MEMS accelerometer on I2C bus.",
        "Total Acceleration Magnitude:",
        "   Total Accel = √( X² + Y² + Z² )",
        "Dynamic Vibration Deviation:",
        "   Vibration = | Total Accel - 1.0g |",
        "Physical Significance: Normal stationary motor exhibits ~1.0g due to gravity. Mechanical imbalance causes dynamic oscillations around 1.0g."
    ]
    for it in vibe_pts:
        p = tf.add_paragraph()
        p.text = it
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    create_card(s11, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s11.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Vibration Severity Classification"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    classes = [
        ("LOW (< 0.05 g)", "Smooth operation, balanced rotor, healthy bearings.", ACCENT_EMERALD),
        ("MEDIUM (0.05 - 0.15 g)", "Developing minor imbalance, looseness, or slight misalignment. Warning logged.", ACCENT_AMBER),
        ("HIGH (≥ 0.15 g)", "Severe vibration! Bearing raceway damage, rotor eccentricity, or imminent structural failure.", RGBColor(239, 68, 68))
    ]
    for tag, desc, col in classes:
        p = tf.add_paragraph()
        p.text = f"• {tag}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col
        p = tf.add_paragraph()
        p.text = f"  {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s11, "MPU6050 measures 3D acceleration. We calculate vector magnitude minus 1g gravity to isolate vibration amplitude into Low, Med, and High severity.")

    # SLIDE 12: IR SENSOR & L298N MOTOR DRIVER
    s12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s12)
    add_header(s12, 11, "IR Optical Sensing & L298N Motor Driver")
    
    create_card(s12, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s12.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "IR Optical Sensor (GPIO 35)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    ir_pts = [
        "Infrared Transmitter/Receiver Pair: Detects reflective object presence and optical passage.",
        "RPM / Shaft Speed Detection: Counts optical pulses per revolution when paired with an encoder wheel.",
        "Position Gating: Confirms physical presence and rotation of conveyor or motor shaft.",
        "Interfacing: Connected to ESP32 input GPIO 35 with interrupt-driven pulse counter."
    ]
    for it in ir_pts:
        p = tf.add_paragraph()
        p.text = f"• {it}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    create_card(s12, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s12.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "L298N Dual H-Bridge Driver"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    l298_pts = [
        "Power Stage: Dual full-bridge driver handling up to 46V and 2A continuous motor loads.",
        "ESP32 Control Connections:",
        "  • ENA (Speed / PWM): GPIO 25 (0–255 duty cycle)",
        "  • IN1 (Direction A): GPIO 26",
        "  • IN2 (Direction B): GPIO 27",
        "Safety Interlock: Allows backend and edge logic to immediately stop motor upon critical alarm."
    ]
    for it in l298_pts:
        p = tf.add_paragraph()
        p.text = f"• {it}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s12, "IR sensor monitors rotation and presence. L298N drives the DC motor with PWM speed control from GPIO 25, 26, and 27.")

    # SLIDE 13: CIRCUIT DIAGRAM & PIN MAP
    s13 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s13)
    add_header(s13, 12, "Hardware Circuit Schematics & Pin Mapping")
    
    create_card(s13, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tb = s13.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.4), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ESP32 GPIO Pin Configuration"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    pinouts = [
        ("DHT22 DATA", "GPIO 4", "Digital 1-Wire Thermal Bus"),
        ("ACS712 VOUT", "GPIO 34", "Analog ADC1_CH6 (Current Sensor)"),
        ("IR SENSOR OUT", "GPIO 35", "Digital ADC1_CH7 / Input"),
        ("MPU6050 SDA", "GPIO 21", "I2C Data Bus"),
        ("MPU6050 SCL", "GPIO 22", "I2C Clock Bus"),
        ("L298N ENA", "GPIO 25", "Hardware PWM Speed Control"),
        ("L298N IN1", "GPIO 26", "Motor Direction Control Phase A"),
        ("L298N IN2", "GPIO 27", "Motor Direction Control Phase B")
    ]
    for sig, pin, note in pinouts:
        p = tf.add_paragraph()
        p.text = f"• {sig:14} ➔ {pin:8} ({note})"
        p.font.size = Pt(11)
        p.font.name = "Consolas"
        p.font.color.rgb = TEXT_WHITE

    create_card(s13, Inches(7.1), Inches(1.8), Inches(5.4), Inches(5.0))
    tb = s13.shapes.add_textbox(Inches(7.4), Inches(2.0), Inches(4.8), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Electrical Design & Protection"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AMBER
    elec_pts = [
        "Common Ground: Unified GND reference between 12V motor supply, 5V sensor rail, and ESP32 3.3V logic.",
        "Decoupling Capacitors: 100uF + 0.1uF across motor rail to snub inductive kickback.",
        "I2C Pull-Up Resistors: 4.7kΩ on SDA/SCL lines to guarantee clean 400kHz bus communication.",
        "Optical Isolation: ACS712 galvanic isolation prevents inductive voltage spikes from damaging the ESP32."
    ]
    for it in elec_pts:
        p = tf.add_paragraph()
        p.text = f"✔ {it}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s13, "Clear pin map and electrical protection details ensuring stable 24/7 continuous sensor readings.")

    # SLIDE 14: SYSTEM ARCHITECTURE
    s14 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s14)
    add_header(s14, 13, "End-to-End System Architecture")
    
    blocks = [
        ("1. PHYSICAL SENSORS", ["DC Industrial Motor", "DHT22 / ACS712", "MPU6050 / IR Sensor", "L298N Motor Driver"], ACCENT_CYAN),
        ("2. EDGE COMPUTING", ["Sensor Acquisition", "Offset Calibration", "Wi-Fi TCP/IP Stack", "HTTP JSON Payload"], ACCENT_BLUE),
        ("3. FASTAPI BACKEND", ["REST API Ingestion", "WebSocket Broadcaster", "Data Validation", "SQLite Time-Series"], ACCENT_EMERALD),
        ("4. WEB DASHBOARD", ["Live Real-Time Gauges", "Historical Trend Charts", "Alert Notification Bar", "Speed / Power Controls"], ACCENT_AMBER),
        ("5. AI PREDICTION", ["Vibration Anomaly Grading", "Time-Series Degradation", "Failure Risk Probability", "Maintenance Alerts"], RGBColor(168, 85, 247))
    ]
    for idx, (title, items, col) in enumerate(blocks):
        left = Inches(0.8 + idx * 2.4)
        top = Inches(1.8)
        create_card(s14, left, top, Inches(2.25), Inches(5.0))
        tb = s14.shapes.add_textbox(left + Inches(0.1), top + Inches(0.2), Inches(2.05), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = col
        for it in items:
            p = tf.add_paragraph()
            p.text = f"• {it}"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s14, "Architecture moves seamlessly from physical motor sensors to ESP32 edge processing, FastAPI cloud backend, React dashboard, and ML prediction.")

    # SLIDE 15: ESP32 EMBEDDED CODING
    s15 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s15)
    add_header(s15, 14, "ESP32 Firmware Execution Lifecycle")
    
    create_card(s15, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s15.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Microcontroller Program Flow"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    steps = [
        "1. Power-on Self Test & Pin Configuration",
        "2. Initialize MPU6050 on I2C (0x68) & DHT22",
        "3. Connect to Wi-Fi SSID with auto-reconnect fallback",
        "4. Non-Blocking Loop (Every 2000 ms):",
        "   • Sample ADC for ACS712 current & compute RMS",
        "   • Read MPU6050 accelerometer X, Y, Z vector",
        "   • Read DHT22 temperature (°C) & humidity (%)",
        "   • Classify vibration level (Low / Med / High)",
        "   • Construct JSON payload & POST to FastAPI backend"
    ]
    for s in steps:
        p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE

    create_card(s15, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s15.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sample Edge JSON Telemetry"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    code_snippet = (
        "{\n"
        '  "motor_id": "MOTOR-01",\n'
        '  "temperature": 34.4,\n'
        '  "humidity": 62.2,\n'
        '  "current": 2.40,\n'
        '  "mpu_x": 0.259,\n'
        '  "mpu_y": -0.965,\n'
        '  "mpu_z": -0.062,\n'
        '  "total_accel": 1.001,\n'
        '  "vibration": 0.038,\n'
        '  "vibration_level": "LOW",\n'
        '  "ir_status": "HIGH",\n'
        '  "motor_state": "ON"\n'
        "}"
    )
    p = tf.add_paragraph()
    p.text = code_snippet
    p.font.size = Pt(11)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(224, 242, 254)
    add_speaker_notes(s15, "Firmware runs a non-blocking 2-second acquisition loop, builds structured JSON, and posts telemetry via HTTP REST.")

    # SLIDE 16: SENSOR DATA PROCESSING PIPELINE
    s16 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s16)
    add_header(s16, 15, "Edge Sensor Data Processing & Filtering")
    
    procs = [
        ("Current Data Pipeline", [
            "1. Read 12-bit ADC raw count (0 - 4095).",
            "2. Convert ADC counts to voltage.",
            "3. Subtract calibrated zero-point (1.8925 V).",
            "4. Divide by sensor sensitivity.",
            "5. Apply multi-sample moving average filter to eliminate motor brush noise."
        ], ACCENT_CYAN),
        ("Vibration Vector Math", [
            "1. Read raw 16-bit registers for X, Y, Z.",
            "2. Scale by LSB sensitivity to g-units.",
            "3. Compute total vector magnitude:\n   √( X² + Y² + Z² )",
            "4. Calculate deviation from 1.0g gravity.",
            "5. Classify severity (Low / Med / High)."
        ], ACCENT_EMERALD),
        ("Threshold Rules Engine", [
            "Temperature Rules:\n • < 50°C : Normal\n • 50-70°C : Warning\n • > 70°C : Critical Alarm",
            "Vibration Rules:\n • < 0.05g : Low\n • 0.05-0.15g : Medium\n • ≥ 0.15g : High Alarm",
            "Current Rules:\n • > 4.5A : Overload Trip"
        ], ACCENT_AMBER)
    ]
    for idx, (title, items, col) in enumerate(procs):
        left = Inches(0.8 + idx * 4.0)
        top = Inches(1.8)
        create_card(s16, left, top, Inches(3.7), Inches(5.0))
        tb = s16.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), Inches(3.4), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        for it in items:
            p = tf.add_paragraph()
            p.text = it
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s16, "Detailed look at mathematical conversion and rule-based thresholding for current, vibration, and temperature.")

    # SLIDE 17: FASTAPI BACKEND SERVER
    s17 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s17)
    add_header(s17, 16, "FastAPI Backend Architecture & Endpoints")
    
    create_card(s17, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s17.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Backend Core Architecture"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    be_pts = [
        "Asynchronous Processing: High-concurrency FastAPI framework handles hundreds of sensor POSTs per second.",
        "Schema Validation: Pydantic models validate incoming sensor ranges and reject malformed packets.",
        "SQLite Persistence: Time-series database records every telemetry reading with microsecond timestamps.",
        "WebSocket Broadcasting: Instantly pushes fresh sensor readings to all connected web dashboard clients."
    ]
    for it in be_pts:
        p = tf.add_paragraph()
        p.text = f"• {it}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    create_card(s17, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s17.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "REST & WebSocket Endpoints"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    endpoints = [
        ("POST /api/telemetry", "Ingests ESP32 sensor readings"),
        ("GET /api/motor/status", "Returns live motor status & telemetry"),
        ("GET /api/motor/history", "Queries historical time-series data"),
        ("POST /api/motor/control", "Dispatches ON/OFF & speed PWM"),
        ("WS /ws/telemetry", "Real-time bi-directional streaming"),
        ("GET /api/predict", "Runs ML inference on latest sensor window")
    ]
    for ep, desc in endpoints:
        p = tf.add_paragraph()
        p.text = f"✔ {ep}\n   ➔ {desc}"
        p.font.size = Pt(11)
        p.font.name = "Consolas"
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s17, "FastAPI backend ingests telemetry, validates with Pydantic, saves to SQLite, and streams over WebSockets.")

    # SLIDE 18: WEB DASHBOARD INTERFACE
    s18 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s18)
    add_header(s18, 17, "Web Monitoring Dashboard")
    
    create_card(s18, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s18.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Live Dashboard UI Snapshot"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    dash_snapshot = (
        "┌──────────────────────────────────────────────┐\n"
        "│  MOTOR-01  |  STATUS: ● ONLINE (RUNNING)    │\n"
        "├──────────────────────────────────────────────┤\n"
        "│  Temperature:  34.4 °C     [ NORMAL ]       │\n"
        "│  Humidity:     62.2 %      [ AMBIENT ]      │\n"
        "│  Current Draw: 2.40 A      [ NOMINAL ]      │\n"
        "│  Vibration:    LOW (0.038g)[ HEALTHY ]      │\n"
        "│  MPU-X: 0.259g | MPU-Y: -0.965g | MPU-Z: -0.062g │\n"
        "│  IR Sensor:    OBJECT DETECTED              │\n"
        "└──────────────────────────────────────────────┘"
    )
    p = tf.add_paragraph()
    p.text = dash_snapshot
    p.font.size = Pt(10)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(224, 242, 254)

    create_card(s18, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s18.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Operator Features & Controls"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    features = [
        "Live Sensor Telemetry: Instant sub-second gauge updates.",
        "Interactive Charts: Time-series graphs for Temperature, Current, and Vibration trends.",
        "Remote Motor Control: Start, Stop, and Speed slider PWM modulation.",
        "Alert Notification Center: Real-time visual toasts and audible alarm thresholds.",
        "Historical Export: Download CSV datasets for external analysis."
    ]
    for it in features:
        p = tf.add_paragraph()
        p.text = f"✔ {it}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s18, "The dashboard provides plant operators with real-time health indicators, live waveforms, and motor control switches.")

    # SLIDE 19: MACHINE LEARNING & PREDICTION
    s19 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s19)
    add_header(s19, 18, "Predictive Maintenance & Machine Learning")
    
    create_card(s19, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s19.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Predictive ML Pipeline"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    ml_steps = [
        "1. Data Collection: Aggregation of multi-sensor time-series.",
        "2. Preprocessing & Normalization: Outlier rejection and baseline normalization.",
        "3. Feature Extraction:",
        "   • Time-domain: RMS, Peak-to-Peak, Crest Factor.",
        "   • Frequency-domain: FFT harmonic peaks.",
        "4. Model Architecture: Random Forest / Isolation Forest / Gradient Boosting.",
        "5. Output Prediction: Remaining Useful Life (RUL) & Failure Probability %."
    ]
    for it in ml_steps:
        p = tf.add_paragraph()
        p.text = it
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    create_card(s19, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb = s19.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Condition States & Failure Modes"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    modes = [
        ("NORMAL STATE", "All sensor metrics within baseline operational limits.", ACCENT_EMERALD),
        ("WARNING STATE", "Elevated vibration (0.05-0.15g) or gradual temperature drift. Maintenance scheduled.", ACCENT_AMBER),
        ("CRITICAL ANOMALY", "High vibration (>0.15g) or current spike (>4.5A). Auto-shutdown triggered.", RGBColor(239, 68, 68))
    ]
    for tag, desc, col in modes:
        p = tf.add_paragraph()
        p.text = f"● {tag}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col
        p = tf.add_paragraph()
        p.text = f"  {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s19, "ML models analyze feature vectors from vibration, current, and temperature to forecast equipment failure before it happens.")

    # SLIDE 20: HARDWARE & EXPERIMENTAL OUTPUT PHOTOS
    s20 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s20)
    add_header(s20, 19, "Hardware Assembly & Experimental Setup")
    
    photo_slots = [
        ("1. Assembled Hardware Node", "ESP32 + DHT22 + ACS712 + MPU6050 integrated breadboard prototype.", ACCENT_CYAN),
        ("2. Motor & Sensor Placement", "Physical sensor coupling on industrial DC motor casing.", ACCENT_BLUE),
        ("3. Arduino Serial Monitor", "Real-time edge terminal showing live calibrated telemetry.", ACCENT_EMERALD),
        ("4. Wi-Fi & Backend Ingestion", "FastAPI server receiving HTTP POST payloads.", ACCENT_AMBER),
        ("5. Live Web Dashboard", "Operator UI displaying real-time gauges and controls.", RGBColor(168, 85, 247)),
        ("6. Predictive Anomaly Output", "ML classification results and condition status badges.", RGBColor(236, 72, 153))
    ]
    for idx, (title, desc, col) in enumerate(photo_slots):
        row = idx // 3
        col_idx = idx % 3
        left = Inches(0.8 + col_idx * 4.0)
        top = Inches(1.8 + row * 2.5)
        create_card(s20, left, top, Inches(3.7), Inches(2.2))
        tb = s20.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(3.3), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col
        p = tf.add_paragraph()
        p.text = f"\n[Insert Actual Photo Here]\n\n{desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s20, "Insert your actual high-resolution hardware photos, serial monitor captures, and dashboard screenshots into these placeholders.")

    # SLIDE 21: SYSTEM OUTPUT & TELEMETRY
    s21 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s21)
    add_header(s21, 20, "Live System Output & Telemetry Log")
    
    create_card(s21, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tb = s21.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.4), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ESP32 Terminal Output Stream"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    serial_log = (
        "[ESP32] Wi-Fi Connected: 192.168.1.104\n"
        "[SENSORS] DHT22: Temp=34.40°C, Hum=62.20%\n"
        "[ACS712] ADC Raw=2348, CalibV=1.8925V, I=2.40A\n"
        "[MPU6050] Accel X=0.259g, Y=-0.965g, Z=-0.062g\n"
        "[MATH] Total Accel = 1.001g, Vibration = 0.038g\n"
        "[SEVERITY] Level = LOW (Healthy)\n"
        "[IR SENSOR] State = HIGH (Object Present)\n"
        "[HTTP] POST /api/telemetry -> 200 OK (38ms)\n"
        "[MOTOR] State = ON, PWM Speed = 65%\n"
        "--------------------------------------------"
    )
    p = tf.add_paragraph()
    p.text = serial_log
    p.font.size = Pt(11)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(224, 242, 254)

    create_card(s21, Inches(7.1), Inches(1.8), Inches(5.4), Inches(5.0))
    tb = s21.shapes.add_textbox(Inches(7.4), Inches(2.0), Inches(4.8), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Telemetry Validation & Performance"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    pts = [
        "Update Frequency: Consistent 2.0 second acquisition cycles with zero missed frames.",
        "Low Latency: Average HTTP POST transmission latency under 50 ms over local Wi-Fi.",
        "Zero Drift: Running average filtering prevents false positive alarms during startup.",
        "Robustness: Automatic Wi-Fi reconnection handling if network dropouts occur."
    ]
    for it in pts:
        p = tf.add_paragraph()
        p.text = f"✔ {it}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s21, "Demonstrates the exact live telemetry stream generated by the ESP32 and successfully processed by the backend.")

    # SLIDE 22: SYSTEM ADVANTAGES
    s22 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s22)
    add_header(s22, 21, "Key Advantages of the IoT Solution")
    
    advs = [
        ("Real-Time 24/7 Monitoring", "Continuous telemetry collection eliminates blind spots between manual inspection intervals.", ACCENT_CYAN),
        ("Early Fault Detection", "Identifies subtle thermal, electrical, and vibrational shifts before catastrophic failure.", ACCENT_BLUE),
        ("Reduced Downtime & Costs", "Enables condition-based maintenance, minimizing costly emergency repairs and lost output.", ACCENT_EMERALD),
        ("Remote Accessibility", "Web dashboard allows plant managers to monitor fleet health from any workstation or mobile device.", ACCENT_AMBER),
        ("Scalable Architecture", "FastAPI backend and SQLite/PostgreSQL architecture easily scales to hundreds of motor nodes.", RGBColor(168, 85, 247)),
        ("Industry 4.0 Ready", "Seamlessly integrates with smart factory MES/SCADA infrastructure and predictive AI systems.", RGBColor(236, 72, 153))
    ]
    for idx, (title, desc, col) in enumerate(advs):
        row = idx // 3
        col_idx = idx % 3
        left = Inches(0.8 + col_idx * 4.0)
        top = Inches(1.8 + row * 2.5)
        create_card(s22, left, top, Inches(3.7), Inches(2.2))
        tb = s22.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(3.3), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"✔ {title}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s22, "Highlighting 24/7 continuous visibility, early fault detection, downtime reduction, and scalability.")

    # SLIDE 23: CONCLUSION
    s23 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s23)
    add_header(s23, 22, "Project Conclusion")
    
    create_card(s23, Inches(0.8), Inches(1.8), Inches(11.733), Inches(3.6))
    tb = s23.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(11.1), Inches(3.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Summary of Achievements"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    concl = [
        "Successfully developed and validated a full-stack IoT industrial motor monitoring platform.",
        "Demonstrated real-time sensor data acquisition (Temperature, Humidity, Current, 3D Acceleration, Optical IR) on the ESP32.",
        "Built a robust, high-performance FastAPI backend with SQLite persistence and WebSocket broadcasting.",
        "Engineered an interactive web dashboard providing plant operators with live health gauges and remote motor controls.",
        "Established the baseline telemetry foundation and feature extraction pipeline for predictive maintenance ML models."
    ]
    for it in concl:
        p = tf.add_paragraph()
        p.text = f"✔ {it}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    create_card(s23, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.4), bg_color=RGBColor(24, 39, 75))
    tb = s23.shapes.add_textbox(Inches(1.1), Inches(5.7), Inches(11.1), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Final Project Statement:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p = tf.add_paragraph()
    p.text = "“The project demonstrates how modern IoT edge computing, cloud backend integration, and machine learning can converge to transform industrial motor monitoring from reactive maintenance into a proactive, condition-based predictive intelligence ecosystem.”"
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = RGBColor(224, 242, 254)
    add_speaker_notes(s23, "Conclude by reiterating how edge IoT and ML come together to deliver real-time condition-based intelligence.")

    # SLIDE 24: FUTURE ENHANCEMENTS
    s24 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s24)
    add_header(s24, 23, "Future Enhancements & Roadmap")
    
    enhancements = [
        ("1. Advanced ML Fault Classification", "Train deep neural networks and gradient boosting on expanded failure datasets to isolate bearing faults, imbalance, and winding breakdown.", ACCENT_CYAN),
        ("2. Multi-Motor Fleet Scalability", "Expand node architecture with unique UUIDs (M001, M002, ...) for plant-wide multi-motor fleet monitoring.", ACCENT_BLUE),
        ("3. Cloud & Edge Deployment", "Deploy backend on AWS/GCP with MQTT brokers and edge AI microcontrollers (ESP32-S3 with TinyML).", ACCENT_EMERALD),
        ("4. Mobile Operator Application", "Develop cross-platform React Native / Flutter apps with push notifications and SMS alert dispatches.", ACCENT_AMBER),
        ("5. Frequency-Domain Vibration (FFT)", "Implement high-frequency onboard FFT vibration spectral analysis to detect specific bearing defect frequencies.", RGBColor(168, 85, 247)),
        ("6. Industrial-Grade Transducers", "Upgrade prototype sensors to DIN-rail industrial 4-20mA and Modbus RTU RS-485 vibration transmitters.", RGBColor(236, 72, 153))
    ]
    for idx, (title, desc, col) in enumerate(enhancements):
        row = idx // 3
        col_idx = idx % 3
        left = Inches(0.8 + col_idx * 4.0)
        top = Inches(1.8 + row * 2.5)
        create_card(s24, left, top, Inches(3.7), Inches(2.2))
        tb = s24.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(3.3), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
    add_speaker_notes(s24, "Roadmap includes TinyML at the edge, FFT vibration spectra, mobile apps, and industrial-grade sensor integration.")

    # SLIDE 25: THANK YOU & Q&A
    s25 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s25)
    
    create_card(s25, Inches(1.8), Inches(1.2), Inches(9.733), Inches(5.1), bg_color=RGBColor(24, 39, 75))
    tb = s25.shapes.add_textbox(Inches(2.2), Inches(1.6), Inches(8.933), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "THANK YOU!"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "IoT-Based Motor Monitoring and Predictive Maintenance System"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nQuestions & Technical Discussion"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nSense  ➔  Process  ➔  Communicate  ➔  Store  ➔  Visualize  ➔  Predict  ➔  Maintain"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(224, 242, 254)
    p.alignment = PP_ALIGN.CENTER
    
    add_speaker_notes(s25, "Thank the audience, panel, and faculty. Open the floor for technical questions regarding hardware, software, or ML.")

    # Output file
    output_path = "IoT_Motor_Monitoring_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {os.path.abspath(output_path)}")

if __name__ == "__main__":
    build_presentation()
